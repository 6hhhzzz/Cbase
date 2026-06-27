"""PPTX 解析器 — 新建，使用 python-pptx 提取幻灯片文本和表格。

借鉴 RAGFlow 的 PPT 解析思路：
    - 每张幻灯片独立为一个组（以标题 TextBlock 分隔）
    - 提取占位符文本（标题、正文）
    - 提取表格为 TableBlock
    - 保留幻灯片序号
"""

from pptx import Presentation

from common import get_logger
from parsing.base import BaseParser
from parsing.models import ParsedDocument, DocumentMetadata, TextBlock, TableBlock

logger = get_logger(__name__)


class PptxParser(BaseParser):
    """PPTX 文件解析器。提取每张幻灯片的文本和表格。"""

    async def _do_parse(self, file_path: str) -> ParsedDocument:
        prs = Presentation(file_path)
        blocks: list[TextBlock | TableBlock] = []
        total_slides = len(prs.slides)

        for slide_idx, slide in enumerate(prs.slides):
            slide_blocks = _extract_slide(slide, slide_idx)
            if slide_blocks:
                # 幻灯片分隔标题
                blocks.append(TextBlock(
                    text=f"--- 幻灯片 {slide_idx + 1}/{total_slides} ---",
                    page_num=slide_idx,
                    layout_type="title",
                    level=1,
                ))
                blocks.extend(slide_blocks)

        # 尝试提取演示文稿标题
        title = None
        if prs.core_properties.title:
            title = prs.core_properties.title

        logger.info(
            f"PPTX 解析完成: {file_path}, slides={total_slides}, blocks={len(blocks)}"
        )
        return ParsedDocument(
            blocks=blocks,
            metadata=DocumentMetadata(
                file_type="pptx",
                file_name=file_path,
                page_count=total_slides,
                title=title,
            ),
        )

    def supports(self, file_type: str) -> bool:
        return file_type == "pptx"


def _extract_slide(slide, slide_idx: int) -> list[TextBlock | TableBlock]:
    """提取单张幻灯片的内容。"""
    blocks: list[TextBlock | TableBlock] = []

    for shape in slide.shapes:
        # ---- 表格 ----
        if shape.has_table:
            table = shape.table
            data = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                if any(c for c in cells):
                    data.append(cells)
            if data:
                html = _to_html_table(data)
                desc = _describe_slide_table(data, slide_idx)
                blocks.append(TableBlock(
                    html=html,
                    description=desc,
                    page_num=slide_idx,
                ))
            continue

        # ---- 文本框 ----
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if not text:
                continue

            # 判断是否为标题
            is_title = shape.is_placeholder and (
                shape.placeholder_format.type == 1  # TITLE
            )
            layout_type = "title" if is_title else "text"
            level = 1 if is_title else None

            blocks.append(TextBlock(
                text=text,
                page_num=slide_idx,
                layout_type=layout_type,
                level=level,
            ))

    return blocks


def _to_html_table(data: list[list[str]]) -> str:
    """二维数组 → HTML <table>。"""
    if not data:
        return ""
    rows = ["<table>"]
    rows.append("<thead><tr>")
    for cell in data[0]:
        rows.append(f"<th>{cell}</th>")
    rows.append("</tr></thead>")
    if len(data) > 1:
        rows.append("<tbody>")
        for row in data[1:]:
            rows.append("<tr>")
            for cell in row:
                rows.append(f"<td>{cell}</td>")
            rows.append("</tr>")
        rows.append("</tbody>")
    rows.append("</table>")
    return "\n".join(rows)


def _describe_slide_table(data: list[list[str]], slide_idx: int) -> str:
    """生成幻灯片表格语义描述。"""
    headers = data[0] if data else []
    row_count = len(data) - 1
    col_count = len(headers)
    return (
        f"幻灯片 {slide_idx + 1} 中的表格：{row_count} 行 × {col_count} 列"
        f"，表头 [{', '.join(headers[:5])}]"
    )
